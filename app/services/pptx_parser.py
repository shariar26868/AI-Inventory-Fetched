import logging
from typing import List, Dict, Any

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    """Parse a .pptx file or fail clearly when dependency is missing."""
    if Presentation is None:
        raise RuntimeError(
            "Missing dependency python-pptx. Install it with `pip install python-pptx`."
        )
    """
    Extract text content from PowerPoint slides. Returns one entry per slide
    with `raw_text` and `_slide` metadata.
    """
    try:
        prs = Presentation(file_path)
        results: List[Dict[str, Any]] = []

        for i, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if parts:
                results.append({
                    "raw_text": "\n".join(parts),
                    "_slide": i,
                    "_source": "pptx",
                })

        logger.info(f"PPTX parsed: {len(results)} slides with text")
        return results

    except Exception as e:
        logger.error(f"PPTX parse error: {e}")
        raise
